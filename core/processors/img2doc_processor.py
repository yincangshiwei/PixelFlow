"""
图片转文档处理器 —— 图片排版导出 (PPT / PDF / Word)
支持：
  - 导出为 PPTX / PDF / DOCX
  - 自定义文档尺寸（如 33.87 x 19.05 cm）
  - 单页图片数量（1/2/4 张）及排版方式
  - 排序方式：默认排序 / Excel 指定列排序 / 分组排序
  - 图片压缩（按目标大小，二分法，先压缩再排版）
  - 叠加层系统（每页可叠加文字/图片，支持自定义坐标）
"""
import os
import io
import math
import copy
import json
import traceback
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QGroupBox,
    QLabel, QComboBox, QSpinBox, QDoubleSpinBox,
    QPushButton, QFileDialog, QLineEdit, QCheckBox,
    QFrame, QColorDialog, QSizePolicy, QListWidget,
    QListWidgetItem, QStackedWidget
)
from PySide6.QtCore import Qt, Signal
from PySide6.QtGui import QColor

from core.base_processor import BaseProcessor, ProcessResult, register_processor
from core.image_processor import compress_to_target_size
import config


def _import_pptx():
    import pptx
    from pptx.util import Cm, Pt
    return pptx, Cm, Pt

def _import_docx():
    import docx
    from docx.shared import Cm, Pt
    return docx, Cm, Pt

def _import_reportlab():
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import cm
    return canvas, cm

def _import_openpyxl():
    import openpyxl
    return openpyxl


def _col2idx(col_str):
    """将 Excel 列字母转为 0-based 索引，如 A->0, B->1, AA->26"""
    num = 0
    for c in col_str:
        num = num * 26 + (ord(c.upper()) - ord('A')) + 1
    return num - 1


class OverlayLayer:
    """叠加层基类"""
    def __init__(self, layer_type, x_cm, y_cm, placement='overlay', name=''):
        self.layer_type = layer_type   # 'text' or 'image'
        self.x_cm = x_cm              # X坐标（cm，相对页面左上角）
        self.y_cm = y_cm              # Y坐标（cm，相对页面左上角）
        self.placement = placement    # 'overlay'=叠加（不占空间）/ 'reserved'=占用空间（参与排版边距计算）
        self.name = name              # 用户自定义名称（列表显示|后面的文字）


class TextLayer(OverlayLayer):
    """文字叠加层"""
    def __init__(self, x_cm=1.0, y_cm=17.0, source='fixed', text='',
                 excel_file='', excel_col='C', match_column=1, data_column=2,
                 excel_row_start=2, font_family='Microsoft YaHei', font_size_pt=12,
                 bold=False, color='#000000', placement='overlay', name=''):
        super().__init__('text', x_cm, y_cm, placement, name)
        self.source = source           # 'fixed', 'excel', 'filename'
        self.text = text               # 固定文本
        self.excel_file = excel_file   # Excel文件路径
        self.excel_col = excel_col     # Excel列字母（保留兼容）
        self.match_column = match_column
        self.data_column = data_column
        self.excel_row_start = excel_row_start
        self.font_family = font_family
        self.font_size_pt = font_size_pt
        self.bold = bold
        self.color = color


class ImageLayer(OverlayLayer):
    """图片叠加层"""
    def __init__(self, x_cm=26.0, y_cm=0.5, path='', w_cm=5.0, h_cm=3.0, placement='overlay', name=''):
        super().__init__('image', x_cm, y_cm, placement, name)
        self.path = path
        self.w_cm = w_cm   # 宽度（cm）
        self.h_cm = h_cm   # 高度（cm）


class LayerListItem(QListWidgetItem):
    """叠加层列表项"""
    def __init__(self, layer, index):
        super().__init__()
        self.layer = layer
        self.index = index
        self.update_text()

    def update_text(self):
        label = self.layer.name if self.layer.name else ""
        if self.layer.layer_type == 'text':
            self.setText(f"📝 文字 {self.index + 1} | {label}")
        else:
            self.setText(f"🖼️ 图片 {self.index + 1} | {label}")


@register_processor
class Img2DocProcessor(BaseProcessor):
    name = "图片排版导出"
    description = "将多张图片按指定布局排版，导出为 PPT / PDF / Word"
    icon = "📑"
    preset_id = "img2doc"

    @property
    def is_batch_processor(self) -> bool:
        return True

    def __init__(self):
        self._panel: QWidget | None = None
        self._layers: list[OverlayLayer] = []
        self._list_widget: QListWidget | None = None
        self._config_stack: QStackedWidget | None = None

    def create_panel(self, parent=None) -> QWidget:
        self._panel = QWidget(parent)
        root = QVBoxLayout(self._panel)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(8)

        # ══ 1. 图片压缩（先压缩再排版，可选，默认不选）══
        grp_compress = QGroupBox("图片压缩（排版前预处理）")
        c_lay = QVBoxLayout(grp_compress)

        c_row1 = QHBoxLayout()
        self.chk_compress = QCheckBox("启用压缩（按目标大小，先压缩再排版）")
        self.chk_compress.setChecked(False)
        c_row1.addWidget(self.chk_compress)
        c_row1.addStretch()
        c_lay.addLayout(c_row1)

        self.widget_compress = QWidget()
        c_inner = QHBoxLayout(self.widget_compress)
        c_inner.setContentsMargins(0, 0, 0, 0)
        c_inner.addWidget(QLabel("目标大小(KB):"))
        self.spin_target_kb = QSpinBox()
        self.spin_target_kb.setRange(10, 102400)
        self.spin_target_kb.setValue(500)
        self.spin_target_kb.setSingleStep(50)
        self.spin_target_kb.setToolTip("每张图片压缩后的目标大小上限（KB），使用二分法寻找最优质量")
        c_inner.addWidget(self.spin_target_kb)
        c_inner.addWidget(QLabel("压缩格式:"))
        self.combo_compress_fmt = QComboBox()
        self.combo_compress_fmt.addItems(["JPEG", "WEBP"])
        self.combo_compress_fmt.setToolTip("PNG 不支持有损压缩，建议选 JPEG 或 WEBP")
        self.combo_compress_fmt.setStyleSheet(config.COMBOBOX_STYLE)
        c_inner.addWidget(self.combo_compress_fmt)
        c_inner.addStretch()
        c_lay.addWidget(self.widget_compress)

        def _on_compress_toggled(checked):
            self.widget_compress.setVisible(checked)
        self.chk_compress.toggled.connect(_on_compress_toggled)
        _on_compress_toggled(False)

        root.addWidget(grp_compress)

        # ══ 2. 输出格式与尺寸 ══
        grp_fmt = QGroupBox("输出设置")
        f_lay = QVBoxLayout(grp_fmt)

        row1 = QHBoxLayout()
        row1.addWidget(QLabel("导出格式:"))
        self.combo_fmt = QComboBox()
        self.combo_fmt.addItems(["PPTX", "PDF", "DOCX"])
        self.combo_fmt.setStyleSheet(config.COMBOBOX_STYLE)
        row1.addWidget(self.combo_fmt)
        row1.addStretch()
        f_lay.addLayout(row1)

        row2 = QHBoxLayout()
        row2.addWidget(QLabel("页面宽度(cm):"))
        self.spin_w = QSpinBox()
        self.spin_w.setRange(5, 200)
        self.spin_w.setValue(33)
        row2.addWidget(self.spin_w)
        row2.addWidget(QLabel("高度(cm):"))
        self.spin_h = QSpinBox()
        self.spin_h.setRange(5, 200)
        self.spin_h.setValue(19)
        row2.addWidget(self.spin_h)
        row2.addStretch()
        f_lay.addLayout(row2)
        root.addWidget(grp_fmt)

        def _on_fmt_changed(text):
            if text == "PPTX":
                self.spin_w.setValue(33)
                self.spin_h.setValue(19)
            else:
                self.spin_w.setValue(21)
                self.spin_h.setValue(30)
        self.combo_fmt.currentTextChanged.connect(_on_fmt_changed)

        # ══ 3. 排版布局 ══
        grp_layout = QGroupBox("排版布局")
        l_lay = QVBoxLayout(grp_layout)

        l_row1 = QHBoxLayout()
        l_row1.addWidget(QLabel("每页图片数:"))
        self.spin_count = QSpinBox()
        self.spin_count.setRange(1, 100)
        self.spin_count.setValue(1)
        self.spin_count.setToolTip("输入每页要排列的图片数量，将自动计算最佳网格布局")
        l_row1.addWidget(self.spin_count)
        l_row1.addWidget(QLabel("每行最大列数:"))
        self.spin_max_cols = QSpinBox()
        self.spin_max_cols.setRange(1, 20)
        self.spin_max_cols.setValue(4)
        self.spin_max_cols.setToolTip("每行最多放几张图，超出后自动换行")
        l_row1.addWidget(self.spin_max_cols)
        l_row1.addStretch()
        l_lay.addLayout(l_row1)

        l_row2 = QHBoxLayout()
        self.chk_keep_ratio = QCheckBox("保持图片比例 (不拉伸)")
        self.chk_keep_ratio.setChecked(True)
        l_row2.addWidget(self.chk_keep_ratio)
        l_row2.addWidget(QLabel("行内排列方向:"))
        self.combo_row_align = QComboBox()
        self.combo_row_align.addItem("从左到右", "left")
        self.combo_row_align.addItem("居中对齐", "center")
        self.combo_row_align.addItem("从右到左", "right")
        self.combo_row_align.setToolTip(
            "从左到右：图片紧靠左侧排列\n"
            "居中对齐：图片在行内水平居中\n"
            "从右到左：图片紧靠右侧，顺序仍从左到右"
        )
        self.combo_row_align.setStyleSheet(config.COMBOBOX_STYLE)
        l_row2.addWidget(self.combo_row_align)
        l_row2.addStretch()
        l_lay.addLayout(l_row2)

        root.addWidget(grp_layout)

        # ══ 4. 排序规则 ══
        grp_sort = QGroupBox("排序规则")
        s_lay = QVBoxLayout(grp_sort)

        s_row1 = QHBoxLayout()
        s_row1.addWidget(QLabel("排序模式:"))
        self.combo_sort = QComboBox()
        self.combo_sort.addItem("按文件名默认排序", "default")
        self.combo_sort.addItem("按 Excel 指定列排序", "excel_sort")
        self.combo_sort.setStyleSheet(config.COMBOBOX_STYLE)
        s_row1.addWidget(self.combo_sort)
        s_row1.addStretch()
        s_lay.addLayout(s_row1)

        self.widget_excel = QWidget()
        e_lay = QVBoxLayout(self.widget_excel)
        e_lay.setContentsMargins(0, 0, 0, 0)

        e_row1 = QHBoxLayout()
        self.txt_excel = QLineEdit()
        self.txt_excel.setPlaceholderText("选择 Excel 文件 (.xlsx)")
        self.txt_excel.setReadOnly(True)
        e_row1.addWidget(self.txt_excel)
        self.btn_excel = QPushButton("浏览...")
        self.btn_excel.clicked.connect(self._browse_excel)
        e_row1.addWidget(self.btn_excel)
        e_lay.addLayout(e_row1)

        e_row2 = QHBoxLayout()
        e_row2.addWidget(QLabel("匹配文件名列(字母):"))
        self.txt_col_name = QLineEdit("A")
        self.txt_col_name.setMaximumWidth(40)
        e_row2.addWidget(self.txt_col_name)
        e_row2.addWidget(QLabel("排序列(字母):"))
        self.txt_col_val = QLineEdit("B")
        self.txt_col_val.setMaximumWidth(40)
        e_row2.addWidget(self.txt_col_val)
        e_row2.addWidget(QLabel("数据起始行:"))
        self.spin_start_row = QSpinBox()
        self.spin_start_row.setRange(1, 99999)
        self.spin_start_row.setValue(2)
        self.spin_start_row.setMaximumWidth(60)
        self.spin_start_row.setToolTip("填 2 表示第 1 行是表头，从第 2 行开始读取数据")
        e_row2.addWidget(self.spin_start_row)
        e_row2.addStretch()
        e_lay.addLayout(e_row2)
        s_lay.addWidget(self.widget_excel)

        def _on_sort_changed(idx):
            mode = self.combo_sort.currentData()
            self.widget_excel.setVisible(mode != "default")
        self.combo_sort.currentIndexChanged.connect(_on_sort_changed)
        _on_sort_changed(0)
        root.addWidget(grp_sort)

        # ══ 5. 叠加层系统（左右分栏）══
        grp_overlay = QGroupBox("页面叠加元素（每页均叠加）")
        ov_lay = QHBoxLayout(grp_overlay)
        ov_lay.setSpacing(10)

        # ── 左侧：元素列表 + 图标按钮 ──
        left_panel = QWidget()
        left_panel.setMaximumWidth(280)
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(0, 0, 0, 0)
        left_layout.setSpacing(6)

        self._list_widget = QListWidget()
        self._list_widget.setMinimumHeight(140)
        self._list_widget.currentRowChanged.connect(self._on_layer_selected)
        left_layout.addWidget(self._list_widget)

        # 操作按钮（emoji图标）
        btn_layout = QHBoxLayout()
        btn_layout.setSpacing(4)

        def _make_icon_btn(emoji, tooltip, slot):
            btn = QPushButton(emoji)
            btn.setToolTip(tooltip)
            btn.setStyleSheet(config.ICON_BTN_STYLE)
            btn.setCursor(Qt.PointingHandCursor)
            btn.clicked.connect(slot)
            return btn

        btn_add_text = _make_icon_btn("📝", "添加文字", self._add_text_layer)
        btn_layout.addWidget(btn_add_text)

        btn_add_image = _make_icon_btn("🖼", "添加图片", self._add_image_layer)
        btn_layout.addWidget(btn_add_image)

        btn_delete = _make_icon_btn("🗑", "删除选中", self._delete_layer)
        btn_layout.addWidget(btn_delete)

        btn_layout.addSpacing(4)

        btn_rename = _make_icon_btn("✏", "重命名", self._rename_layer)
        btn_layout.addWidget(btn_rename)

        btn_layout.addSpacing(6)

        btn_move_up = _make_icon_btn("⬆", "上移", self._move_layer_up)
        btn_layout.addWidget(btn_move_up)

        btn_move_down = _make_icon_btn("⬇", "下移", self._move_layer_down)
        btn_layout.addWidget(btn_move_down)

        btn_layout.addStretch()
        left_layout.addLayout(btn_layout)

        ov_lay.addWidget(left_panel, stretch=2)

        # ── 右侧：元素配置面板 ──
        grp_config = QGroupBox("元素配置")
        config_root = QVBoxLayout(grp_config)
        config_root.setContentsMargins(8, 8, 8, 8)

        self._config_stack = QStackedWidget()

        # 默认提示页面
        hint_page = QWidget()
        hint_layout = QVBoxLayout(hint_page)
        hint_label = QLabel("请先添加元素，然后在左侧列表中选择一个元素进行配置")
        hint_label.setStyleSheet("color:#666e88;font-size:12px;font-style:italic;")
        hint_label.setAlignment(Qt.AlignCenter)
        hint_layout.addWidget(hint_label)
        self._config_stack.addWidget(hint_page)

        config_root.addWidget(self._config_stack)
        ov_lay.addWidget(grp_config, stretch=5)

        # 坐标说明
        coord_hint = QLabel(
            "坐标说明：X/Y 为距页面左上角的距离（厘米），宽/高为元素尺寸（厘米）。",
            wordWrap=True
        )
        coord_hint.setStyleSheet("color:#666e88;font-size:11px;")
        root.addWidget(grp_overlay)
        root.addWidget(coord_hint)

        root.addStretch()
        return self._panel

    def _browse_excel(self):
        f, _ = QFileDialog.getOpenFileName(self._panel, "选择 Excel", str(Path.home() / "Desktop"), "Excel (*.xlsx)")
        if f:
            self.txt_excel.setText(f)

    def _create_text_config_page(self, layer: TextLayer) -> QWidget:
        """创建文字层配置页面"""
        page = QWidget()
        page_layout = QVBoxLayout(page)
        page_layout.setContentsMargins(0, 0, 0, 0)
        page_layout.setSpacing(8)

        # ── 添加方式 ──
        placement_layout = QHBoxLayout()
        placement_layout.addWidget(QLabel("添加方式:"))
        self._combo_placement = QComboBox()
        self._combo_placement.addItem("叠加（不占空间）", "overlay")
        self._combo_placement.addItem("占用空间（参与排版边距）", "reserved")
        self._combo_placement.setCurrentIndex(self._combo_placement.findData(layer.placement))
        self._combo_placement.setToolTip(
            "叠加：直接叠在图片排版区上方，不影响图片排版区域大小\n"
            "占用空间：元素会占据页面边缘区域，图片排版区自动缩小避让"
        )
        self._combo_placement.setStyleSheet(config.COMBOBOX_STYLE)
        placement_layout.addWidget(self._combo_placement)
        placement_layout.addStretch()
        page_layout.addLayout(placement_layout)

        # ── 数据源 ──
        src_layout = QHBoxLayout()
        src_layout.addWidget(QLabel("数据来源:"))
        combo_source = QComboBox()
        combo_source.addItem("固定文本", "fixed")
        combo_source.addItem("Excel列数据", "excel")
        combo_source.addItem("图片文件名称", "filename")
        combo_source.setCurrentIndex(combo_source.findData(layer.source))
        combo_source.setStyleSheet(config.COMBOBOX_STYLE)
        src_layout.addWidget(combo_source)
        src_layout.addStretch()
        page_layout.addLayout(src_layout)

        # 固定文本内容
        self._widget_fixed_text = QWidget()
        fixed_layout = QVBoxLayout(self._widget_fixed_text)
        fixed_layout.setContentsMargins(0, 0, 0, 0)
        fixed_layout.addWidget(QLabel("文本内容:"))
        self._text_content = QLineEdit(layer.text)
        self._text_content.setPlaceholderText("输入要显示的文字")
        fixed_layout.addWidget(self._text_content)
        page_layout.addWidget(self._widget_fixed_text)

        # Excel配置
        self._widget_excel = QWidget()
        excel_layout = QVBoxLayout(self._widget_excel)
        excel_layout.setContentsMargins(0, 0, 0, 0)

        # Excel文件选择
        file_layout = QHBoxLayout()
        file_layout.addWidget(QLabel("Excel文件:"))
        self._excel_file_input = QLineEdit(layer.excel_file)
        self._excel_file_input.setPlaceholderText("选择Excel文件...")
        file_layout.addWidget(self._excel_file_input)
        btn_browse_excel = QPushButton("浏览")
        btn_browse_excel.clicked.connect(self._browse_excel_for_layer)
        file_layout.addWidget(btn_browse_excel)
        excel_layout.addLayout(file_layout)

        # 匹配列和数据列
        excel_col_layout = QHBoxLayout()
        excel_col_layout.addWidget(QLabel("匹配列:"))
        self._excel_match_column = QSpinBox()
        self._excel_match_column.setRange(1, 100)
        self._excel_match_column.setValue(layer.match_column if layer.match_column > 0 else 1)
        self._excel_match_column.setToolTip("图片文件名所在的列号（从1开始），用于匹配对应行")
        excel_col_layout.addWidget(self._excel_match_column)

        excel_col_layout.addWidget(QLabel("数据列:"))
        self._excel_data_column = QSpinBox()
        self._excel_data_column.setRange(1, 100)
        self._excel_data_column.setValue(layer.data_column if layer.data_column > 0 else 2)
        self._excel_data_column.setToolTip("要读取的文本数据列号（从1开始）")
        excel_col_layout.addWidget(self._excel_data_column)
        excel_col_layout.addStretch()
        excel_layout.addLayout(excel_col_layout)

        # 起始行
        excel_row_layout = QHBoxLayout()
        excel_row_layout.addWidget(QLabel("数据起始行:"))
        self._excel_row_start = QSpinBox()
        self._excel_row_start.setRange(1, 10000)
        self._excel_row_start.setValue(layer.excel_row_start)
        self._excel_row_start.setToolTip("数据从第几行开始（跳过表头）")
        excel_row_layout.addWidget(self._excel_row_start)
        excel_row_layout.addStretch()
        excel_layout.addLayout(excel_row_layout)

        # Excel使用说明
        excel_hint = QLabel("💡 匹配列：图片文件名所在列 | 数据列：要叠加的文本内容所在列")
        excel_hint.setStyleSheet("color:#666e88;font-size:11px;font-style:italic;")
        excel_hint.setWordWrap(True)
        excel_layout.addWidget(excel_hint)

        page_layout.addWidget(self._widget_excel)

        # 文件名提示
        self._widget_filename_hint = QLabel("💡 将使用当前页第一张图片的文件名（不含扩展名）作为文本内容")
        self._widget_filename_hint.setStyleSheet("color:#666e88;font-size:11px;font-style:italic;")
        self._widget_filename_hint.setWordWrap(True)
        page_layout.addWidget(self._widget_filename_hint)

        def _on_source_change(idx):
            source = combo_source.itemData(idx)
            self._widget_fixed_text.setVisible(source == 'fixed')
            self._widget_excel.setVisible(source == 'excel')
            self._widget_filename_hint.setVisible(source == 'filename')
            layer.source = source

        combo_source.currentIndexChanged.connect(_on_source_change)
        _on_source_change(combo_source.currentIndex())

        # 字体设置
        font_layout = QHBoxLayout()
        font_layout.addWidget(QLabel("字体:"))
        self._font_family = QComboBox()
        # 常用字体列表（中文字体优先）
        fonts = ['Microsoft YaHei', 'SimHei', 'SimSun', 'KaiTi', 'FangSong', 'Arial', 'Times New Roman']
        self._font_family.addItems(fonts)
        # 恢复自定义字体
        if hasattr(self, '_custom_fonts'):
            for fn in self._custom_fonts:
                if self._font_family.findText(fn) < 0:
                    self._font_family.addItem(fn)
        idx = self._font_family.findText(layer.font_family)
        if idx >= 0:
            self._font_family.setCurrentIndex(idx)
        else:
            idx = self._font_family.findText('Microsoft YaHei')
            if idx >= 0:
                self._font_family.setCurrentIndex(idx)
        
        self._font_family.setStyleSheet(config.COMBOBOX_STYLE)

        # 添加加载字体按钮
        btn_load_font = QPushButton("加载字体")
        btn_load_font.clicked.connect(self._load_font)
        font_layout.addWidget(btn_load_font)
        font_layout.addWidget(self._font_family)

        font_layout.addWidget(QLabel("大小:"))
        self._font_size_pt = QSpinBox()
        self._font_size_pt.setRange(6, 200)
        self._font_size_pt.setValue(layer.font_size_pt)
        font_layout.addWidget(self._font_size_pt)

        self._chk_bold = QCheckBox("加粗")
        self._chk_bold.setChecked(layer.bold)
        font_layout.addWidget(self._chk_bold)

        font_layout.addStretch()
        page_layout.addLayout(font_layout)

        # 颜色选择
        color_layout = QHBoxLayout()
        color_layout.addWidget(QLabel("文字颜色:"))
        self._text_color_btn = QPushButton()
        self._text_color_btn.setFixedSize(36, 26)
        self._text_color_btn.setCursor(Qt.PointingHandCursor)
        self._text_color_btn.clicked.connect(self._pick_text_color)
        color_layout.addWidget(self._text_color_btn)

        self._text_color_lbl = QLabel(layer.color)
        self._text_color_lbl.setStyleSheet("color:#b0b8c8;font-size:12px;background:transparent;")
        color_layout.addWidget(self._text_color_lbl)
        color_layout.addStretch()
        page_layout.addLayout(color_layout)
        self._refresh_text_color_btn(layer.color)

        # ── 位置（cm）──
        pos_layout = QHBoxLayout()
        pos_layout.addWidget(QLabel("X(cm):"))
        self._text_x_cm = QDoubleSpinBox()
        self._text_x_cm.setRange(0, 200)
        self._text_x_cm.setValue(layer.x_cm)
        self._text_x_cm.setDecimals(2)
        self._text_x_cm.setMaximumWidth(90)
        self._text_x_cm.setToolTip("距页面左边缘的距离（厘米）")
        pos_layout.addWidget(self._text_x_cm)
        pos_layout.addWidget(QLabel("Y(cm):"))
        self._text_y_cm = QDoubleSpinBox()
        self._text_y_cm.setRange(0, 200)
        self._text_y_cm.setValue(layer.y_cm)
        self._text_y_cm.setDecimals(2)
        self._text_y_cm.setMaximumWidth(90)
        self._text_y_cm.setToolTip("距页面上边缘的距离（厘米）")
        pos_layout.addWidget(self._text_y_cm)
        pos_layout.addStretch()
        page_layout.addLayout(pos_layout)

        self._current_layer = layer
        self._combo_source = combo_source
        return page

    def _create_image_config_page(self, layer: ImageLayer) -> QWidget:
        """创建图片层配置页面"""
        page = QWidget()
        page_layout = QVBoxLayout(page)
        page_layout.setContentsMargins(0, 0, 0, 0)
        page_layout.setSpacing(8)

        # ── 添加方式 ──
        placement_layout = QHBoxLayout()
        placement_layout.addWidget(QLabel("添加方式:"))
        self._combo_placement = QComboBox()
        self._combo_placement.addItem("叠加（不占空间）", "overlay")
        self._combo_placement.addItem("占用空间（参与排版边距）", "reserved")
        self._combo_placement.setCurrentIndex(self._combo_placement.findData(layer.placement))
        self._combo_placement.setToolTip(
            "叠加：直接叠在图片排版区上方，不影响图片排版区域大小\n"
            "占用空间：元素会占据页面边缘区域，图片排版区自动缩小避让"
        )
        self._combo_placement.setStyleSheet(config.COMBOBOX_STYLE)
        placement_layout.addWidget(self._combo_placement)
        placement_layout.addStretch()
        page_layout.addLayout(placement_layout)

        # ── 图片文件选择 ──
        file_layout = QHBoxLayout()
        file_layout.addWidget(QLabel("图片文件:"))
        self._image_file_input = QLineEdit(layer.path)
        self._image_file_input.setPlaceholderText("选择要叠加的图片（PNG/JPG）...")
        self._image_file_input.setReadOnly(True)
        file_layout.addWidget(self._image_file_input)
        btn_browse = QPushButton("浏览")
        btn_browse.clicked.connect(self._browse_overlay_image)
        file_layout.addWidget(btn_browse)
        page_layout.addLayout(file_layout)

        # ── 位置（cm）──
        pos_layout = QHBoxLayout()
        pos_layout.addWidget(QLabel("X(cm):"))
        self._image_x_cm = QDoubleSpinBox()
        self._image_x_cm.setRange(0, 200)
        self._image_x_cm.setValue(layer.x_cm)
        self._image_x_cm.setDecimals(2)
        self._image_x_cm.setMaximumWidth(90)
        self._image_x_cm.setToolTip("距页面左边缘的距离（厘米）")
        pos_layout.addWidget(self._image_x_cm)
        pos_layout.addWidget(QLabel("Y(cm):"))
        self._image_y_cm = QDoubleSpinBox()
        self._image_y_cm.setRange(0, 200)
        self._image_y_cm.setValue(layer.y_cm)
        self._image_y_cm.setDecimals(2)
        self._image_y_cm.setMaximumWidth(90)
        self._image_y_cm.setToolTip("距页面上边缘的距离（厘米）")
        pos_layout.addWidget(self._image_y_cm)
        pos_layout.addStretch()
        page_layout.addLayout(pos_layout)

        # ── 大小（cm）──
        size_layout = QHBoxLayout()
        size_layout.addWidget(QLabel("宽(cm):"))
        self._image_w_cm = QDoubleSpinBox()
        self._image_w_cm.setRange(0.1, 200)
        self._image_w_cm.setValue(layer.w_cm)
        self._image_w_cm.setDecimals(2)
        self._image_w_cm.setMaximumWidth(90)
        self._image_w_cm.setToolTip("图片宽度（厘米）")
        size_layout.addWidget(self._image_w_cm)
        size_layout.addWidget(QLabel("高(cm):"))
        self._image_h_cm = QDoubleSpinBox()
        self._image_h_cm.setRange(0.1, 200)
        self._image_h_cm.setValue(layer.h_cm)
        self._image_h_cm.setDecimals(2)
        self._image_h_cm.setMaximumWidth(90)
        self._image_h_cm.setToolTip("图片高度（厘米）")
        size_layout.addWidget(self._image_h_cm)
        size_layout.addStretch()
        page_layout.addLayout(size_layout)

        self._current_layer = layer
        return page

    def _browse_overlay_image(self):
        """浏览选择叠加图片"""
        file_path, _ = QFileDialog.getOpenFileName(
            self._panel, "选择图片文件", str(Path.home() / "Desktop"),
            "Image Files (*.png *.jpg *.jpeg *.webp)"
        )
        if file_path:
            self._image_file_input.setText(file_path)

    def _browse_excel_for_layer(self):
        """浏览选择Excel文件（用于文字层配置）"""
        file_path, _ = QFileDialog.getOpenFileName(
            self._panel, "选择Excel文件", str(Path.home() / "Desktop"), "Excel Files (*.xlsx *.xls)"
        )
        if file_path:
            self._excel_file_input.setText(file_path)

    def _load_font(self):
        """加载外部字体文件"""
        file_path, _ = QFileDialog.getOpenFileName(
            self._panel, "选择字体文件", str(Path.home() / "Desktop"),
            "Font Files (*.ttf *.otf *.ttc)"
        )
        if not file_path:
            return

        try:
            # 验证字体文件是否有效
            from PIL import ImageFont
            test_font = ImageFont.truetype(file_path, 12)

            # 获取字体名称
            font_name = Path(file_path).stem

            # 检查是否已存在
            existing_fonts = [self._font_family.itemText(i) for i in range(self._font_family.count())]
            if font_name in existing_fonts:
                from PySide6.QtWidgets import QMessageBox
                QMessageBox.warning(self._panel, "提示", f"字体 '{font_name}' 已存在")
                return

            # 添加到字体列表
            self._font_family.addItem(font_name)
            self._font_family.setCurrentIndex(self._font_family.count() - 1)

            # 保存字体路径
            if not hasattr(self, '_custom_fonts'):
                self._custom_fonts = {}
            self._custom_fonts[font_name] = file_path

            from PySide6.QtWidgets import QMessageBox
            QMessageBox.information(self._panel, "成功", f"字体 '{font_name}' 加载成功！")

        except Exception as e:
            from PySide6.QtWidgets import QMessageBox
            QMessageBox.warning(self._panel, "错误", f"字体加载失败: {e}")

    def _pick_text_color(self):
        """选择文本颜色"""
        current_color = self._text_color_lbl.text()
        c = QColorDialog.getColor(QColor(current_color), self._panel, "选择文字颜色")
        if c.isValid():
            color = c.name()
            self._refresh_text_color_btn(color)
            self._text_color_lbl.setText(color.upper())

    def _refresh_text_color_btn(self, color):
        """刷新颜色按钮样式"""
        self._text_color_btn.setStyleSheet(
            f"QPushButton{{background:{color};border:2px solid #5a5a6a;border-radius:6px;min-width:36px;min-height:24px;}}"
            f"QPushButton:hover{{border-color:#5b8af5;}}"
        )

    def _collect_current_layer_options(self):
        """收集当前编辑元素的参数"""
        if not hasattr(self, '_current_layer'):
            return

        layer = self._current_layer

        # 通用：添加方式
        if hasattr(self, '_combo_placement'):
            layer.placement = self._combo_placement.currentData()

        if layer.layer_type == 'text':
            if hasattr(self, '_combo_source'):
                layer.source = self._combo_source.currentData()
            if layer.source == 'fixed' and hasattr(self, '_text_content'):
                layer.text = self._text_content.text()
            elif layer.source == 'excel':
                if hasattr(self, '_excel_file_input'):
                    layer.excel_file = self._excel_file_input.text()
                if hasattr(self, '_excel_match_column'):
                    layer.match_column = self._excel_match_column.value()
                if hasattr(self, '_excel_data_column'):
                    layer.data_column = self._excel_data_column.value()
                if hasattr(self, '_excel_row_start'):
                    layer.excel_row_start = self._excel_row_start.value()
            if hasattr(self, '_font_family'):
                layer.font_family = self._font_family.currentText()
            if hasattr(self, '_font_size_pt'):
                layer.font_size_pt = self._font_size_pt.value()
            if hasattr(self, '_chk_bold'):
                layer.bold = self._chk_bold.isChecked()
            if hasattr(self, '_text_color_lbl'):
                layer.color = self._text_color_lbl.text()
            if hasattr(self, '_text_x_cm'):
                layer.x_cm = self._text_x_cm.value()
            if hasattr(self, '_text_y_cm'):
                layer.y_cm = self._text_y_cm.value()
        else:
            if hasattr(self, '_image_file_input'):
                layer.path = self._image_file_input.text()
            if hasattr(self, '_image_x_cm'):
                layer.x_cm = self._image_x_cm.value()
            if hasattr(self, '_image_y_cm'):
                layer.y_cm = self._image_y_cm.value()
            if hasattr(self, '_image_w_cm'):
                layer.w_cm = self._image_w_cm.value()
            if hasattr(self, '_image_h_cm'):
                layer.h_cm = self._image_h_cm.value()

    # ── 叠加层 UI 管理 ────────────────────────────────────────────────────────

    def _add_text_layer(self):
        """添加一个文字叠加层"""
        text_count = sum(1 for e in self._layers if e.layer_type == 'text')
        default_y = 17.0 - (text_count * 1.5)  # 依次向上排列，避免重叠
        layer = TextLayer(x_cm=1.0, y_cm=max(0.5, default_y))
        self._layers.append(layer)
        item = LayerListItem(layer, len(self._layers) - 1)
        self._list_widget.addItem(item)
        self._list_widget.setCurrentRow(self._list_widget.count() - 1)

    def _add_image_layer(self):
        """添加一个图片叠加层"""
        img_count = sum(1 for e in self._layers if e.layer_type == 'image')
        default_x = 26.0 - (img_count * 6.0)  # 依次向左排列
        layer = ImageLayer(x_cm=max(0.5, default_x), y_cm=0.5)
        self._layers.append(layer)
        item = LayerListItem(layer, len(self._layers) - 1)
        self._list_widget.addItem(item)
        self._list_widget.setCurrentRow(self._list_widget.count() - 1)

    def _delete_layer(self):
        """删除选中的元素"""
        row = self._list_widget.currentRow()
        if row < 0:
            return

        # 临时断开信号，避免 takeItem 触发 currentRowChanged
        self._list_widget.currentRowChanged.disconnect(self._on_layer_selected)

        self._list_widget.takeItem(row)
        self._layers.pop(row)

        # 更新索引
        for i in range(self._list_widget.count()):
            item = self._list_widget.item(i)
            item.index = i
            item.update_text()

        # 清除配置面板
        self._clear_config_pages()

        # 重新连接信号
        self._list_widget.currentRowChanged.connect(self._on_layer_selected)

        # 如果还有元素，自动选中并显示配置
        if self._list_widget.count() > 0:
            new_row = min(row, self._list_widget.count() - 1)
            # 先强制设为 -1，确保 setCurrentRow 一定能触发 currentRowChanged
            self._list_widget.setCurrentRow(-1)
            self._list_widget.setCurrentRow(new_row)
        else:
            self._config_stack.setCurrentIndex(0)

    def _rename_layer(self):
        """重命名选中元素的名称（| 后面的文字）"""
        row = self._list_widget.currentRow()
        if row < 0:
            return
        from PySide6.QtWidgets import QInputDialog
        current_name = self._layers[row].name
        new_name, ok = QInputDialog.getText(
            self._panel, "重命名元素", "请输入新名称:",
            text=current_name
        )
        if ok:
            self._layers[row].name = new_name.strip()
            self._list_widget.item(row).update_text()

    def _move_layer_up(self):
        """上移选中的元素"""
        row = self._list_widget.currentRow()
        if row <= 0:
            return

        # 临时断开信号
        self._list_widget.currentRowChanged.disconnect(self._on_layer_selected)

        item = self._list_widget.takeItem(row)
        self._list_widget.insertItem(row - 1, item)
        self._layers[row], self._layers[row - 1] = self._layers[row - 1], self._layers[row]

        for i in range(self._list_widget.count()):
            self._list_widget.item(i).index = i
            self._list_widget.item(i).update_text()

        # 重新连接信号
        self._list_widget.currentRowChanged.connect(self._on_layer_selected)

        # 保持选中状态
        self._list_widget.setCurrentRow(row - 1)

    def _move_layer_down(self):
        """下移选中的元素"""
        row = self._list_widget.currentRow()
        if row < 0 or row >= self._list_widget.count() - 1:
            return

        # 临时断开信号
        self._list_widget.currentRowChanged.disconnect(self._on_layer_selected)

        item = self._list_widget.takeItem(row)
        self._list_widget.insertItem(row + 1, item)
        self._layers[row], self._layers[row + 1] = self._layers[row + 1], self._layers[row]

        for i in range(self._list_widget.count()):
            self._list_widget.item(i).index = i
            self._list_widget.item(i).update_text()

        # 重新连接信号
        self._list_widget.currentRowChanged.connect(self._on_layer_selected)

        # 保持选中状态
        self._list_widget.setCurrentRow(row + 1)

    def _on_layer_selected(self, row):
        """选中元素时显示配置面板"""
        if row < 0:
            return

        # 先收集当前编辑元素的参数
        self._collect_current_layer_options()

        layer = self._layers[row]

        # 清除旧配置页面
        self._clear_config_pages()

        # 创建新的配置页面
        if layer.layer_type == 'text':
            config_page = self._create_text_config_page(layer)
        else:
            config_page = self._create_image_config_page(layer)

        self._config_stack.addWidget(config_page)
        self._config_stack.setCurrentIndex(1)

    def _clear_config_pages(self):
        """彻底清除所有配置页面（保留索引0的提示页面）"""
        attrs_to_clear = [
            '_current_layer', '_combo_source', '_combo_placement',
            '_widget_fixed_text', '_widget_excel', '_widget_filename_hint',
            '_text_content', '_excel_file_input', '_excel_match_column',
            '_excel_data_column', '_excel_row_start', '_font_family',
            '_font_size_pt', '_chk_bold', '_text_color_btn', '_text_color_lbl',
            '_text_x_cm', '_text_y_cm',
            '_image_file_input', '_image_x_cm', '_image_y_cm',
            '_image_w_cm', '_image_h_cm'
        ]
        for attr in attrs_to_clear:
            if hasattr(self, attr):
                delattr(self, attr)

        while self._config_stack.count() > 1:
            widget = self._config_stack.widget(1)
            self._config_stack.removeWidget(widget)
            widget.setParent(None)
            widget.deleteLater()

    def process(self, img, options):
        # 此处理器为批量处理器，不会调用单张 process()
        pass

    def gather_options(self) -> dict:
        # 先收集当前正在编辑的元素
        self._collect_current_layer_options()

        # 序列化所有叠加层
        layers_data = []
        for layer in self._layers:
            if layer.layer_type == 'text':
                layers_data.append({
                    'type': 'text',
                    'name': layer.name,
                    'placement': layer.placement,
                    'source': layer.source,
                    'text': layer.text,
                    'excel_file': layer.excel_file,
                    'excel_col': layer.excel_col,
                    'match_column': layer.match_column,
                    'data_column': layer.data_column,
                    'excel_row_start': layer.excel_row_start,
                    'font_family': layer.font_family,
                    'font_size_pt': layer.font_size_pt,
                    'bold': layer.bold,
                    'color': layer.color,
                    'x_cm': layer.x_cm,
                    'y_cm': layer.y_cm,
                })
            else:
                layers_data.append({
                    'type': 'image',
                    'name': layer.name,
                    'placement': layer.placement,
                    'path': layer.path,
                    'x_cm': layer.x_cm,
                    'y_cm': layer.y_cm,
                    'w_cm': layer.w_cm,
                    'h_cm': layer.h_cm,
                })

        # 收集自定义字体信息
        custom_fonts = {}
        if hasattr(self, '_custom_fonts'):
            custom_fonts = self._custom_fonts.copy()

        return {
            "format": self.combo_fmt.currentText(),
            "width_cm": self.spin_w.value(),
            "height_cm": self.spin_h.value(),
            "count_per_page": self.spin_count.value(),
            "max_cols": self.spin_max_cols.value(),
            "keep_ratio": self.chk_keep_ratio.isChecked(),
            "row_align": self.combo_row_align.currentData(),
            "sort_mode": self.combo_sort.currentData(),
            "excel_path": self.txt_excel.text(),
            "col_name": self.txt_col_name.text().upper(),
            "col_val": self.txt_col_val.text().upper(),
            "start_row": self.spin_start_row.value(),
            # 压缩
            "compress_enabled": self.chk_compress.isChecked(),
            "compress_target_kb": self.spin_target_kb.value(),
            "compress_format": self.combo_compress_fmt.currentText(),
            # 叠加层
            "overlay_layers": layers_data,
            "custom_fonts": custom_fonts,
        }

    def default_options(self) -> dict:
        return {
            "format": "PPTX",
            "width_cm": 33,
            "height_cm": 19,
            "count_per_page": 1,
            "max_cols": 4,
            "keep_ratio": True,
            "row_align": "left",
            "sort_mode": "default",
            "excel_path": "",
            "col_name": "A",
            "col_val": "B",
            "start_row": 2,
            "compress_enabled": False,
            "compress_target_kb": 500,
            "compress_format": "JPEG",
            "overlay_layers": [],
            "custom_fonts": {},
        }

    def apply_options(self, options: dict):
        if self._panel is None: return
        fmt = options.get("format", "PPTX")
        idx = self.combo_fmt.findText(fmt)
        if idx >= 0: self.combo_fmt.setCurrentIndex(idx)

        self.spin_w.setValue(options.get("width_cm", 33))
        self.spin_h.setValue(options.get("height_cm", 19))
        self.spin_count.setValue(options.get("count_per_page", 1))
        self.spin_max_cols.setValue(options.get("max_cols", 4))
        self.chk_keep_ratio.setChecked(options.get("keep_ratio", True))
        align_idx = self.combo_row_align.findData(options.get("row_align", "left"))
        if align_idx >= 0:
            self.combo_row_align.setCurrentIndex(align_idx)

        mode = options.get("sort_mode", "default")
        midx = self.combo_sort.findData(mode)
        if midx >= 0: self.combo_sort.setCurrentIndex(midx)

        self.txt_excel.setText(options.get("excel_path", ""))
        self.txt_col_name.setText(options.get("col_name", "A"))
        self.txt_col_val.setText(options.get("col_val", "B"))
        self.spin_start_row.setValue(options.get("start_row", 2))

        # 压缩
        self.chk_compress.setChecked(options.get("compress_enabled", False))
        self.spin_target_kb.setValue(options.get("compress_target_kb", 500))
        ci = self.combo_compress_fmt.findText(options.get("compress_format", "JPEG"))
        if ci >= 0: self.combo_compress_fmt.setCurrentIndex(ci)

        # 叠加层
        # 先收集当前编辑元素的参数
        self._collect_current_layer_options()

        # 清除配置页面，显示提示页面
        self._clear_config_pages()
        self._config_stack.setCurrentIndex(0)

        # 清除现有元素
        self._layers.clear()
        self._list_widget.clear()

        # 加载元素
        layers_data = options.get("overlay_layers", [])
        for data in layers_data:
            if data['type'] == 'text':
                layer = TextLayer(
                    x_cm=data.get('x_cm', 1.0),
                    y_cm=data.get('y_cm', 17.0),
                    source=data.get('source', 'fixed'),
                    text=data.get('text', ''),
                    name=data.get('name', ''),
                    excel_file=data.get('excel_file', ''),
                    excel_col=data.get('excel_col', 'C'),
                    match_column=data.get('match_column', 1),
                    data_column=data.get('data_column', 2),
                    excel_row_start=data.get('excel_row_start', 2),
                    font_family=data.get('font_family', 'Microsoft YaHei'),
                    font_size_pt=data.get('font_size_pt', 12),
                    bold=data.get('bold', False),
                    color=data.get('color', '#000000'),
                    placement=data.get('placement', 'overlay'),
                )
            else:
                layer = ImageLayer(
                    x_cm=data.get('x_cm', 26.0),
                    y_cm=data.get('y_cm', 0.5),
                    path=data.get('path', ''),
                    name=data.get('name', ''),
                    w_cm=data.get('w_cm', 5.0),
                    h_cm=data.get('h_cm', 3.0),
                    placement=data.get('placement', 'overlay'),
                )

            self._layers.append(layer)
            item = LayerListItem(layer, len(self._layers) - 1)
            self._list_widget.addItem(item)

        # 恢复自定义字体
        custom_fonts = options.get('custom_fonts', {})
        if custom_fonts:
            self._custom_fonts = custom_fonts

    def get_output_format(self) -> str:
        return ""  # 批量处理器由其自身决定输出后缀

    def process_batch(self, file_list: list[str], options: dict, output_dir: str, progress_cb=None) -> list[ProcessResult]:
        if not file_list:
            return []

        # 1. 整理文件列表（排序）
        sort_mode = options.get("sort_mode", "default")
        groups = {}

        if sort_mode == "default":
            groups["排版导出"] = sorted(file_list, key=lambda x: Path(x).name)
        else:
            excel_path = options.get("excel_path", "")
            if not os.path.exists(excel_path):
                raise ValueError(f"Excel 文件不存在: {excel_path}")

            openpyxl = _import_openpyxl()
            wb = openpyxl.load_workbook(excel_path, data_only=True, read_only=True)
            try:
                ws = wb.active

                c_name = _col2idx(options.get("col_name", "A"))
                c_val = _col2idx(options.get("col_val", "B"))
                start_row = options.get("start_row", 2)

                mapping = {}
                for row_idx, row in enumerate(ws.iter_rows(values_only=True), start=1):
                    if row_idx < start_row:
                        continue
                    if row[c_name] is not None:
                        mapping[str(row[c_name]).strip()] = row[c_val]
            finally:
                wb.close()

            def get_sort_val(fpath):
                name = Path(fpath).stem
                return mapping.get(name, mapping.get(Path(fpath).name, ""))

            if sort_mode == "excel_sort":
                sorted_files = sorted(file_list, key=lambda x: str(get_sort_val(x)))
                groups["排版导出"] = sorted_files

        # 2. 读取叠加层 Excel 数据（如果有文字层需要从 Excel 读取）
        overlay_layers = options.get("overlay_layers", [])
        # 构建 excel 叠加层的数据映射
        # overlay_excel_data[layer_index] = {col_letter: [val_row2, val_row3, ...]}
        overlay_excel_data = {}
        excel_layers = [(i, l) for i, l in enumerate(overlay_layers) if l.get("type") == "text" and l.get("source") == "excel"]

        for layer_idx, layer in excel_layers:
            excel_file = layer.get("excel_file", "")
            if not excel_file or not os.path.exists(excel_file):
                # 如果没有配置Excel文件，尝试使用排序Excel
                excel_file = options.get("excel_path", "")

            if excel_file and os.path.exists(excel_file):
                openpyxl = _import_openpyxl()
                wb2 = openpyxl.load_workbook(excel_file, data_only=True, read_only=True)
                try:
                    ws2 = wb2.active

                    match_col = layer.get("match_column", 1)
                    data_col = layer.get("data_column", 2)
                    start_row2 = layer.get("excel_row_start", 2)

                    # 构建文件名到数据列的映射
                    data_map = {}
                    for row_idx, row in enumerate(ws2.iter_rows(values_only=True), start=1):
                        if row_idx < start_row2:
                            continue
                        match_val = row[match_col - 1] if match_col - 1 < len(row) else None
                        data_val = row[data_col - 1] if data_col - 1 < len(row) else None
                        if match_val is not None:
                            match_str = str(match_val).strip()
                            # 去除扩展名
                            if '.' in match_str:
                                match_str = Path(match_str).stem
                            data_map[match_str.lower()] = str(data_val) if data_val is not None else ""
                finally:
                    wb2.close()

                overlay_excel_data[layer_idx] = data_map

        # 3. 图片压缩预处理
        compress_enabled = options.get("compress_enabled", False)
        compress_target_kb = options.get("compress_target_kb", 500)
        compress_fmt = options.get("compress_format", "JPEG")
        # 将压缩后的临时路径映射存起来
        _compressed_cache = {}  # original_path -> compressed_bytes (io.BytesIO)

        if compress_enabled:
            all_files = list({f for grp in groups.values() for f in grp})
            for fi, fpath in enumerate(all_files):
                if progress_cb:
                    progress_cb(fi, len(all_files), f"压缩中: {Path(fpath).name}")
                try:
                    with Image.open(fpath) as img:
                        if compress_fmt.upper() in ("JPEG", "JPG"):
                            img_conv = img.convert("RGB")
                            fmt_save = "JPEG"
                        else:
                            img_conv = img.convert("RGBA") if img.mode == "RGBA" else img.convert("RGB")
                            fmt_save = "WEBP"
                        _, quality, _ = compress_to_target_size(img_conv, compress_target_kb, fmt_save)
                        buf = io.BytesIO()
                        img_conv.save(buf, format=fmt_save, quality=quality)
                        buf.seek(0)
                        _compressed_cache[fpath] = buf
                except Exception as e:
                    if progress_cb:
                        progress_cb(fi, len(all_files), f"压缩失败，使用原图: {Path(fpath).name} ({e})")

        # 4. 导出处理
        fmt = options.get("format", "PPTX")
        results = []

        # 进度以「文件总数」为单位，让进度条与文件列表数量对应
        total_files = len(file_list)
        processed_files = 0

        for group_name, files in groups.items():
            count_per_page = options.get("count_per_page", 1)
            total_pages = math.ceil(len(files) / count_per_page)

            if progress_cb:
                progress_cb(processed_files, total_files,
                            f"开始导出: {group_name}（共 {len(files)} 张图，{total_pages} 页）")

            safe_name = "".join(c for c in group_name if c not in r'\/:*?"<>|')
            ext = f".{fmt.lower()}"
            out_path = Path(output_dir) / f"{safe_name}{ext}"

            counter = 1
            while out_path.exists():
                out_path = Path(output_dir) / f"{safe_name}_{counter}{ext}"
                counter += 1

            try:
                # 构造带逐页进度上报的回调
                _group_base = processed_files  # 本组起始已处理数

                def _page_progress_cb(page_idx, page_count, batch_files, _base=_group_base):
                    """每页处理完毕后上报进度和日志"""
                    if progress_cb is None:
                        return
                    imgs_done = min(_base + (page_idx + 1) * options.get("count_per_page", 1), total_files)
                    img_start = _base + page_idx * options.get("count_per_page", 1) + 1
                    img_end = min(img_start + options.get("count_per_page", 1) - 1, total_files)
                    names = ", ".join(Path(f).name for f in batch_files)
                    progress_cb(
                        imgs_done, total_files,
                        f"第 {page_idx + 1}/{page_count} 页  [{img_start}~{img_end}]  {names}"
                    )

                export_kwargs = dict(
                    files=files,
                    out_path=str(out_path),
                    options=options,
                    compressed_cache=_compressed_cache,
                    overlay_layers=overlay_layers,
                    overlay_excel_data=overlay_excel_data,
                    page_progress_cb=_page_progress_cb,
                )
                if fmt == "PPTX":
                    self._export_pptx(**export_kwargs)
                elif fmt == "PDF":
                    self._export_pdf(**export_kwargs)
                elif fmt == "DOCX":
                    self._export_docx(**export_kwargs)

                processed_files += len(files)
                results.append(ProcessResult(
                    input_path=f"分组: {group_name}",
                    output_path=str(out_path),
                    success=True,
                    details={"files_count": len(files), "pages": total_pages}
                ))
            except Exception as e:
                processed_files += len(files)
                results.append(ProcessResult(
                    input_path=f"分组: {group_name}",
                    success=False,
                    error=f"{e}\n{traceback.format_exc()}"
                ))

        if progress_cb:
            progress_cb(total_files, total_files, "完成")

        # 清理压缩缓存，释放 BytesIO 内存
        for buf in _compressed_cache.values():
            buf.close()
        _compressed_cache.clear()

        return results

    def _calc_overlay_margins(self, w_cm, h_cm, overlay_layers):
        """根据 placement='reserved' 的叠加层计算图片排版区域需要避让的边距 (cm)。

        只有标记为"占用空间"的层才参与边距计算，直接使用 cm 坐标。
        判断规则：
        - 文字层：y_cm < h_cm*0.3 视为顶部，y_cm > h_cm*0.7 视为底部
        - 图片层：按 y_cm + h_cm 和 x_cm + w_cm 判断占用区域
        """
        m_top = 0.5    # 基础边距
        m_bottom = 0.5
        m_left = 0.5
        m_right = 0.5

        for layer in overlay_layers:
            if layer.get("placement", "overlay") != "reserved":
                continue  # 叠加模式不参与边距计算

            x_cm_val = layer.get("x_cm", 0)
            y_cm_val = layer.get("y_cm", 0)

            if layer.get("type") == "text":
                # 文字层：估算占用高度 ≈ font_size_pt / 28.35 cm + 0.3cm 间距
                fs = layer.get("font_size_pt", 12)
                text_h_cm = fs / 28.35 + 0.3

                if y_cm_val < h_cm * 0.3:
                    # 顶部区域
                    m_top = max(m_top, y_cm_val + text_h_cm + 0.2)
                elif y_cm_val > h_cm * 0.7:
                    # 底部区域
                    m_bottom = max(m_bottom, h_cm - y_cm_val + 0.2)

            elif layer.get("type") == "image":
                ow_cm = layer.get("w_cm", 5)
                oh_cm = layer.get("h_cm", 3)
                bottom_edge = y_cm_val + oh_cm
                right_edge = x_cm_val + ow_cm

                if bottom_edge <= h_cm * 0.3:
                    # 顶部区域
                    m_top = max(m_top, bottom_edge + 0.2)
                elif y_cm_val >= h_cm * 0.7:
                    # 底部区域
                    m_bottom = max(m_bottom, h_cm - y_cm_val + 0.2)

                if right_edge >= w_cm * 0.7 and bottom_edge <= h_cm * 0.3:
                    # 右上角
                    m_top = max(m_top, bottom_edge + 0.2)
                elif x_cm_val <= w_cm * 0.3 and bottom_edge <= h_cm * 0.3:
                    # 左上角
                    m_top = max(m_top, bottom_edge + 0.2)

        return m_top, m_bottom, m_left, m_right

    def _calc_layout(self, w_cm, h_cm, img_sizes, options=None):
        """计算图片在页面上的精确坐标和尺寸。

        参数:
            w_cm, h_cm: 页面尺寸 (cm)
            img_sizes: [(img_w_px, img_h_px), ...] 每张图的像素尺寸
            options: 包含 max_cols, row_align, overlay_layers 等

        排版规则：
            - 自动根据叠加层位置计算安全边距，图片排版只在安全区域内
            - 宽度按列均分，高度按每张图实际比例计算
            - 逐行计算行高，整体垂直居中
            - 超高时等比缩小
            - row_align: 'left'(从左到右) / 'center'(居中) / 'right'(从右到左)

        返回: (layouts, rows, cols)
        """
        if options is None:
            options = {}

        overlay_layers = options.get("overlay_layers", [])
        m_top, m_bottom, m_left, m_right = self._calc_overlay_margins(w_cm, h_cm, overlay_layers)

        max_cols = options.get("max_cols", 4)
        row_align = options.get("row_align", "left")
        gap = 0.3

        count = len(img_sizes)
        if count == 0:
            return [], 0, 0

        aw = w_cm - m_left - m_right
        ah = h_cm - m_top - m_bottom
        if aw <= 0 or ah <= 0:
            return [], 0, 0

        cols = min(count, max_cols)
        rows = math.ceil(count / cols)

        cell_w = (aw - gap * (cols - 1)) / cols

        # 按行分组，计算每行行高
        row_groups = []
        row_heights = []
        for r in range(rows):
            start = r * cols
            end = min(start + cols, count)
            row_groups.append((start, end))

            max_h = 0
            for idx in range(start, end):
                iw, ih = img_sizes[idx]
                if iw <= 0:
                    iw = 1
                img_h_cm = cell_w * ih / iw
                if img_h_cm > max_h:
                    max_h = img_h_cm
            row_heights.append(max_h)

        total_h = sum(row_heights) + gap * (rows - 1)

        # 超高等比缩小
        scale = 1.0
        if total_h > ah:
            scale = ah / total_h
            cell_w *= scale
            row_heights = [h * scale for h in row_heights]
            total_h = ah

        # 垂直居中
        v_offset = (ah - total_h) / 2

        layouts = []
        current_y = m_top + v_offset
        for r, (start, end) in enumerate(row_groups):
            row_h = row_heights[r]
            row_count = end - start  # 本行实际图片数（最后行可能不满）

            # 计算本行实际总宽度（用于居中/右对齐）
            row_total_w = row_count * cell_w + (row_count - 1) * gap

            if row_align == "center":
                row_x_start = m_left + (aw - row_total_w) / 2
            elif row_align == "right":
                row_x_start = m_left + (aw - row_total_w)
            else:  # left（默认）
                row_x_start = m_left

            for idx in range(start, end):
                c = idx - start
                iw, ih = img_sizes[idx]
                if iw <= 0:
                    iw = 1
                img_w_cm = cell_w
                img_h_cm = cell_w * ih / iw

                x = row_x_start + c * (cell_w + gap)
                y = current_y + (row_h - img_h_cm) / 2
                layouts.append((x, y, img_w_cm, img_h_cm))

            current_y += row_h + gap

        return layouts, rows, cols

    def _get_fit_size(self, img_w, img_h, box_w, box_h):
        """保持比例缩放，返回 (w, h)"""
        scale = min(box_w / img_w, box_h / img_h)
        return img_w * scale, img_h * scale

    # ── 叠加层辅助 ────────────────────────────────────────────────────────────

    def _get_overlay_text(self, layer: dict, page_idx: int, overlay_excel_data: dict, batch_files: list = None, layer_idx: int = 0) -> str:
        """获取文字层在当前页的文字内容

        Args:
            layer: 叠加层配置字典
            page_idx: 页索引
            overlay_excel_data: Excel数据映射 {layer_idx: {filename_lower: data_value}}
            batch_files: 当前页的文件列表（用于filename数据源和Excel匹配）
            layer_idx: 当前层在overlay_layers中的索引
        """
        source = layer.get("source", "fixed")

        if source == "excel":
            # 使用新的Excel数据结构
            data_map = overlay_excel_data.get(layer_idx, {})
            if batch_files:
                # 使用当前页第一张图片的文件名进行匹配
                file_stem = Path(batch_files[0]).stem.lower()
                return data_map.get(file_stem, "")
            return ""
        elif source == "filename":
            # 使用当前页第一张图片的文件名（不含扩展名）
            if batch_files:
                return Path(batch_files[0]).stem
            return ""
        else:  # fixed
            return layer.get("text", "")

    def _open_image_for_export(self, fpath: str, compressed_cache: dict):
        """打开图片，优先使用压缩缓存。返回独立的 Image 对象（不持有文件句柄）"""
        try:
            if fpath in compressed_cache:
                buf = compressed_cache[fpath]
                buf.seek(0)
                img = Image.open(buf)
                img.load()  # 强制读入内存
                return img
            img = Image.open(fpath)
            img.load()
            return img
        except Exception:
            # 兜底：返回一个 1x1 的占位图
            return Image.new("RGB", (1, 1), (128, 128, 128))

    # ── 导出方法 ──────────────────────────────────────────────────────────────

    def _export_pptx(self, files, out_path, options,
                     compressed_cache=None, overlay_layers=None, overlay_excel_data=None,
                     page_progress_cb=None):
        compressed_cache = compressed_cache or {}
        overlay_layers = overlay_layers or []
        overlay_excel_data = overlay_excel_data or {}

        pptx, Cm, Pt = _import_pptx()
        from pptx.util import Pt as PtU
        from pptx.dml.color import RGBColor
        prs = pptx.Presentation()

        w_cm = options["width_cm"]
        h_cm = options["height_cm"]
        prs.slide_width = Cm(w_cm)
        prs.slide_height = Cm(h_cm)
        blank_layout = prs.slide_layouts[6]

        count = options["count_per_page"]
        keep_ratio = options["keep_ratio"]
        total_pages = math.ceil(len(files) / count)

        page_idx = 0
        for i in range(0, len(files), count):
            slide = prs.slides.add_slide(blank_layout)
            batch = files[i:i+count]

            # 读取每张图的实际尺寸
            img_sizes = []
            for f in batch:
                img = self._open_image_for_export(f, compressed_cache)
                img_sizes.append(img.size)
                img.close()

            boxes, _, _ = self._calc_layout(w_cm, h_cm, img_sizes, options)
            for j, f in enumerate(batch):
                bx, by, bw, bh = boxes[j]
                # layout 已按实际比例计算精确尺寸，直接使用
                if f in compressed_cache:
                    compressed_cache[f].seek(0)
                    slide.shapes.add_picture(compressed_cache[f], Cm(bx), Cm(by), Cm(bw), Cm(bh))
                else:
                    slide.shapes.add_picture(f, Cm(bx), Cm(by), Cm(bw), Cm(bh))

            # 叠加层
            for layer_idx, layer in enumerate(overlay_layers):
                x_cm = layer.get("x_cm", 0)
                y_cm = layer.get("y_cm", 0)
                if layer["type"] == "text":
                    text = self._get_overlay_text(layer, page_idx, overlay_excel_data, batch, layer_idx)
                    if not text:
                        continue
                    txBox = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm * 0.5), Cm(h_cm * 0.1))
                    tf = txBox.text_frame
                    tf.word_wrap = False
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = text
                    run.font.size = PtU(layer.get("font_size_pt", 12))
                    color_hex = layer.get("color", "#000000").lstrip("#")
                    run.font.color.rgb = RGBColor(
                        int(color_hex[0:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16)
                    )
                elif layer["type"] == "image":
                    img_path = layer.get("path", "")
                    if not img_path or not os.path.exists(img_path):
                        continue
                    ow_cm = layer.get("w_cm", 5)
                    oh_cm = layer.get("h_cm", 3)
                    slide.shapes.add_picture(img_path, Cm(x_cm), Cm(y_cm), Cm(ow_cm), Cm(oh_cm))

            if page_progress_cb:
                page_progress_cb(page_idx, total_pages, batch)
            page_idx += 1
        prs.save(out_path)

    def _export_pdf(self, files, out_path, options,
                    compressed_cache=None, overlay_layers=None, overlay_excel_data=None,
                    page_progress_cb=None):
        compressed_cache = compressed_cache or {}
        overlay_layers = overlay_layers or []
        overlay_excel_data = overlay_excel_data or {}

        canvas_mod, cm = _import_reportlab()
        from reportlab.lib.colors import HexColor

        w_cm = options["width_cm"]
        h_cm = options["height_cm"]
        c = canvas_mod.Canvas(out_path, pagesize=(w_cm * cm, h_cm * cm))

        count = options["count_per_page"]
        keep_ratio = options["keep_ratio"]
        total_pages = math.ceil(len(files) / count)

        page_idx = 0
        for i in range(0, len(files), count):
            batch = files[i:i+count]

            # 读取每张图的实际尺寸
            img_sizes = []
            for f in batch:
                img = self._open_image_for_export(f, compressed_cache)
                img_sizes.append(img.size)
                img.close()

            boxes, _, _ = self._calc_layout(w_cm, h_cm, img_sizes, options)
            for j, f in enumerate(batch):
                bx, by, bw, bh = boxes[j]
                # PDF 坐标原点在左下角，y 需要翻转
                pdf_y = h_cm - by - bh
                if f in compressed_cache:
                    compressed_cache[f].seek(0)
                    from reportlab.lib.utils import ImageReader
                    reader = ImageReader(compressed_cache[f])
                    try:
                        c.drawImage(reader, bx * cm, pdf_y * cm, width=bw * cm, height=bh * cm)
                    finally:
                        image = getattr(reader, "_image", None)
                        if image is not None:
                            image.close()
                else:
                    c.drawImage(f, bx * cm, pdf_y * cm, width=bw * cm, height=bh * cm)

            # 叠加层
            for layer_idx, layer in enumerate(overlay_layers):
                x_cm_val = layer.get("x_cm", 0)
                y_cm_val = layer.get("y_cm", 0)
                if layer["type"] == "text":
                    text = self._get_overlay_text(layer, page_idx, overlay_excel_data, batch, layer_idx)
                    if not text:
                        continue
                    # PDF 坐标原点在左下，y 需要翻转
                    pdf_ty = (h_cm - y_cm_val) * cm
                    color_hex = layer.get("color", "#000000")
                    c.setFillColor(HexColor(color_hex))
                    font_size = layer.get("font_size_pt", 12)
                    c.setFont("Helvetica", font_size)
                    c.drawString(x_cm_val * cm, pdf_ty, text)
                elif layer["type"] == "image":
                    img_path = layer.get("path", "")
                    if not img_path or not os.path.exists(img_path):
                        continue
                    ow_cm = layer.get("w_cm", 5)
                    oh_cm = layer.get("h_cm", 3)
                    pdf_iy = (h_cm - y_cm_val - oh_cm) * cm
                    c.drawImage(img_path, x_cm_val * cm, pdf_iy, width=ow_cm * cm, height=oh_cm * cm, preserveAspectRatio=True, mask='auto')

            c.showPage()
            if page_progress_cb:
                page_progress_cb(page_idx, total_pages, batch)
            page_idx += 1
        c.save()

    def _export_docx(self, files, out_path, options,
                     compressed_cache=None, overlay_layers=None, overlay_excel_data=None,
                     page_progress_cb=None):
        """DOCX 导出（注：DOCX 不支持绝对坐标叠加，文字层以页脚/段落形式追加）"""
        compressed_cache = compressed_cache or {}
        overlay_layers = overlay_layers or []
        overlay_excel_data = overlay_excel_data or {}

        docx, Cm, Pt = _import_docx()
        doc = docx.Document()

        section = doc.sections[-1]
        section.page_width = Cm(options["width_cm"])
        section.page_height = Cm(options["height_cm"])
        section.top_margin = Cm(0.5)
        section.bottom_margin = Cm(0.5)
        section.left_margin = Cm(0.5)
        section.right_margin = Cm(0.5)

        count = options["count_per_page"]
        keep_ratio = options["keep_ratio"]
        total_pages = math.ceil(len(files) / count)

        page_idx = 0
        for i in range(0, len(files), count):
            batch = files[i:i+count]

            # 读取每张图的实际尺寸
            img_sizes = []
            for f in batch:
                img = self._open_image_for_export(f, compressed_cache)
                img_sizes.append(img.size)
                img.close()

            boxes, rows, cols = self._calc_layout(options["width_cm"], options["height_cm"], img_sizes, options)
            if len(batch) == 1:
                p = doc.add_paragraph()
                p.alignment = 1
                r = p.add_run()
                bx, by, bw, bh = boxes[0]
                if batch[0] in compressed_cache:
                    compressed_cache[batch[0]].seek(0)
                    r.add_picture(compressed_cache[batch[0]], width=Cm(bw), height=Cm(bh))
                else:
                    r.add_picture(batch[0], width=Cm(bw), height=Cm(bh))
            else:
                table = doc.add_table(rows=rows, cols=cols)
                for j, f in enumerate(batch):
                    row_idx = j // cols
                    col_idx = j % cols
                    cell = table.cell(row_idx, col_idx)
                    p = cell.paragraphs[0]
                    p.alignment = 1
                    r = p.add_run()
                    bx, by, bw, bh = boxes[j]
                    if f in compressed_cache:
                        compressed_cache[f].seek(0)
                        r.add_picture(compressed_cache[f], width=Cm(bw), height=Cm(bh))
                    else:
                        r.add_picture(f, width=Cm(bw), height=Cm(bh))

            # DOCX 叠加层：以段落形式追加（DOCX 不支持绝对定位，以注释形式说明位置）
            text_overlays = [(i, l) for i, l in enumerate(overlay_layers) if l.get("type") == "text"]
            img_overlays = [l for l in overlay_layers if l.get("type") == "image"]
            for layer_idx, layer in text_overlays:
                text = self._get_overlay_text(layer, page_idx, overlay_excel_data, batch, layer_idx)
                if text:
                    op = doc.add_paragraph(f"[叠加文字 X:{layer.get('x_cm', 0)}cm Y:{layer.get('y_cm', 0)}cm] {text}")
                    op.runs[0].font.size = Pt(layer.get("font_size_pt", 12))
            for layer in img_overlays:
                img_path = layer.get("path", "")
                if img_path and os.path.exists(img_path):
                    op = doc.add_paragraph(f"[叠加图片 X:{layer.get('x_cm', 0)}cm Y:{layer.get('y_cm', 0)}cm]")
                    op.add_run().add_picture(
                        img_path,
                        width=Cm(layer.get("w_cm", 5))
                    )

            if i + count < len(files):
                doc.add_page_break()
            if page_progress_cb:
                page_progress_cb(page_idx, total_pages, batch)
            page_idx += 1

        doc.save(out_path)